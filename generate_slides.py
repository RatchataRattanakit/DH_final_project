#!/usr/bin/env python3
"""
generate_slides.py  —  DH Project 5-minute presentation
White theme + native python-pptx charts matching index.html data
pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE)

# ── Palette (WHITE THEME) ─────────────────────────────────────────────────
BG      = RGBColor(0xff, 0xff, 0xff)
SURFACE = RGBColor(0xf5, 0xf5, 0xf5)
BORDER  = RGBColor(0xdd, 0xdd, 0xdd)
TEXT    = RGBColor(0x1a, 0x1a, 0x1a)
MUTED   = RGBColor(0x77, 0x77, 0x77)
RED     = RGBColor(0xc4, 0x1e, 0x3a)
GOLD    = RGBColor(0xf0, 0xa5, 0x00)
TEAL    = RGBColor(0x4e, 0xcd, 0xc4)
NAVY    = RGBColor(0x3a, 0x3a, 0x7a)
# Keyword colors (matching index.html)
C_PAT   = RGBColor(0xff, 0x6b, 0x6b)
C_CD    = RGBColor(0xf4, 0xb9, 0x42)
C_NAT   = RGBColor(0x4e, 0xcd, 0xc4)
C_GC    = RGBColor(0x88, 0xcc, 0xaa)
C_BOY   = RGBColor(0xff, 0x8b, 0x94)
C_FOR   = RGBColor(0xcc, 0x33, 0x44)   # foreign brand (darker red for white bg)
C_DOM   = RGBColor(0xd4, 0x92, 0x00)   # domestic brand (darker gold for white bg)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Data (identical to index.html) ────────────────────────────────────────
YEARS        = [2015, 2016, 2017, 2018, 2019, 2020, 2021, 2022, 2023, 2024]
PATRIOTISM   = (246, 283, 352, 405, 502, 601, 686, 755, 745, 759)
CHINA_DREAM  = (312, 400, 483, 543, 640, 736, 805, 881, 867, 880)
NATIONALISM  = (189, 215, 262, 309, 375, 471, 561, 582, 570, 590)
GUOCHAO      = (  8,  11,  16,  44,  97, 239, 496, 418, 449, 445)
BOYCOTT      = ( 28,  36, 109,  36, 134,  57, 329,  79,  61,  68)
FOREIGN_TONE = (-0.82,-1.14,-2.21,-1.58,-2.89,-2.12,-3.94,-2.76,-2.43,-2.62)
DOMESTIC_TONE= ( 0.61, 0.93, 1.42, 1.87, 2.14, 2.58, 3.21, 2.94, 3.12, 3.31)

# era averages (2015-2018 vs 2019-2024)
EARLY = (322, 435, 244,  20,  52)
LATE  = (675, 802, 525, 357, 121)

BOYCOTT_EVENTS = [(2017, "THAAD"), (2019, "NBA"), (2021, "H&M"), (2022, "Olympics")]

# Chart frame dimensions — used to estimate event marker x-positions
CHART_L = 0.45
CHART_T = 1.88
CHART_W = 12.43
CHART_H = 4.6
# Approximate inner plot area (after y-axis label space)
INNER_L = CHART_L + 0.85
INNER_W = CHART_W - 1.1

# ── Helpers ───────────────────────────────────────────────────────────────

def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, text, l, t, w, h,
        size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text, r.font.size, r.font.bold = text, Pt(size), bold
    r.font.italic, r.font.color.rgb  = italic, color
    return tb

def mbox(slide, lines, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for d in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = d.get("align", PP_ALIGN.LEFT)
        if "space_before" in d:
            p.space_before = Pt(d["space_before"])
        r = p.add_run()
        r.text           = d["text"]
        r.font.size      = Pt(d.get("size", 16))
        r.font.bold      = d.get("bold", False)
        r.font.italic    = d.get("italic", False)
        r.font.color.rgb = d.get("color", TEXT)

def rect(slide, l, t, w, h, color=RED, border=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def section_label(slide, text, l=0.5, t=0.28, color=RED):
    box(slide, text, l, t, 12, 0.38, size=10, bold=True, color=color)

def chart_subtitle(slide, text):
    box(slide, text, CHART_L, 1.45, CHART_W, 0.38, size=11, italic=True, color=MUTED)

def finding_callout(slide, text):
    rect(slide, CHART_L, 6.65, CHART_W, 0.72, SURFACE, border=BORDER)
    rect(slide, CHART_L, 6.65, 0.06, 0.72, GOLD)
    box(slide, text, CHART_L + 0.18, 6.73, CHART_W - 0.25, 0.58, size=12, color=TEXT)

def event_markers(slide, events, boycott_only=False):
    """Draw thin gold vertical lines + labels at event years."""
    for year, lbl in events:
        if boycott_only and lbl == "Olympics":
            continue
        frac = (year - 2015) / 9
        ex = INNER_L + frac * INNER_W
        rect(slide, ex - 0.007, CHART_T, 0.014, CHART_H, GOLD)
        box(slide, lbl, ex + 0.05, CHART_T + 0.05, 1.2, 0.3, size=9, bold=True, color=GOLD)

def style_line_chart(chart_obj, series_colors):
    chart_obj.has_title = False

    for i, series in enumerate(chart_obj.series):
        clr = series_colors[i % len(series_colors)]
        series.format.line.color.rgb = clr
        series.format.line.width = Pt(2.2)
        try:
            series.marker.style = XL_MARKER_STYLE.CIRCLE
            series.marker.size  = 5
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = clr
            series.marker.format.line.color.rgb      = clr
        except Exception:
            pass

    try:
        va = chart_obj.value_axis
        va.major_gridlines.format.line.color.rgb = BORDER
        va.tick_labels.font.size      = Pt(10)
        va.tick_labels.font.color.rgb = MUTED
    except Exception:
        pass
    try:
        ca = chart_obj.category_axis
        ca.tick_labels.font.size      = Pt(10)
        ca.tick_labels.font.color.rgb = MUTED
    except Exception:
        pass
    try:
        chart_obj.legend.position          = XL_LEGEND_POSITION.BOTTOM
        chart_obj.legend.include_in_layout = False
    except Exception:
        pass

def style_bar_chart(chart_obj, series_colors):
    chart_obj.has_title = False

    for i, series in enumerate(chart_obj.series):
        clr = series_colors[i % len(series_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = clr
        series.format.line.fill.background()

    try:
        va = chart_obj.value_axis
        va.major_gridlines.format.line.color.rgb = BORDER
        va.tick_labels.font.size      = Pt(10)
        va.tick_labels.font.color.rgb = MUTED
    except Exception:
        pass
    try:
        ca = chart_obj.category_axis
        ca.tick_labels.font.size      = Pt(10)
        ca.tick_labels.font.color.rgb = MUTED
    except Exception:
        pass
    try:
        chart_obj.legend.position          = XL_LEGEND_POSITION.BOTTOM
        chart_obj.legend.include_in_layout = False
    except Exception:
        pass

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1)
rect(s1, 0, 0, 13.33, 0.1, RED)
rect(s1, 0, 7.4, 13.33, 0.1, SURFACE)

section_label(s1, "DIGITAL HUMANITIES  ·  COMPUTATIONAL ANALYSIS  ·  2025", 0.7, 0.38)

mbox(s1, [
    {"text": "How Xi Jinping Embeds Nationalism",  "size": 38, "bold": True, "color": TEXT},
    {"text": "into Chinese Citizens",              "size": 38, "bold": True, "color": TEXT},
    {"text": "Through State Media",                "size": 38, "bold": True, "color": RED, "space_before": 2},
], 0.7, 0.95, 11.8, 3.5)

rect(s1, 0.7, 4.1, 3.0, 0.06, RED)
box(s1, "A computational study of nationalist language in Chinese state media, 2015–2024",
    0.7, 4.32, 11, 0.5, size=15, italic=True, color=MUTED)
box(s1, "Data: GDELT BigQuery  ·  Tools: Python  ·  BigQuery SQL  ·  D3.js",
    0.7, 5.08, 11, 0.38, size=11, color=BORDER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Research Question + Thesis
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
rect(s2, 0, 0, 13.33, 0.1, RED)

section_label(s2, "THE QUESTION")
box(s2,
    "How did Chinese state media change the frequency,\n"
    "context, and emotional framing of nationalist language\n"
    "between 2015 and 2024?",
    0.5, 0.78, 12.3, 2.6, size=27, bold=True, color=TEXT)

rect(s2, 0.5, 3.6, 12.3, 0.04, BORDER)
section_label(s2, "CENTRAL ARGUMENT", 0.5, 3.78, GOLD)

rect(s2, 0.5, 4.22, 12.3, 2.85, SURFACE, border=BORDER)
rect(s2, 0.5, 4.22, 0.07, 2.85, RED)
box(s2,
    "Xi Jinping did not merely inherit Chinese nationalism — he rewired it.\n\n"
    "Through state media, the Party fused patriotic identity with consumer\n"
    "behavior, geopolitical grievance, and everyday life.\n\n"
    "This project measures that fusion computationally.",
    0.72, 4.38, 11.9, 2.6, size=16, italic=True, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data & Method
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
rect(s3, 0, 0, 13.33, 0.1, RED)

section_label(s3, "DATA & METHOD")
box(s3, "What we measured and how", 0.5, 0.68, 9, 0.55, size=23, bold=True, color=TEXT)

rect(s3, 0.45, 1.42, 6.1, 5.75, SURFACE, border=BORDER)
section_label(s3, "DATA SOURCE", 0.72, 1.58, TEAL)
mbox(s3, [
    {"text": "GDELT BigQuery",                             "size": 18, "bold": True, "color": TEXT},
    {"text": "gdelt-bq.gdeltv2.gkg_partitioned",          "size": 11, "color": MUTED, "space_before": 2},
    {"text": " ", "size": 7},
    {"text": "▸  Xinhua · People's Daily · China Daily",  "size": 13, "color": TEXT},
    {"text": "▸  Global Times · CGTN · .cn domains",     "size": 13, "color": TEXT},
    {"text": "▸  2015–2024  (10 years)",                 "size": 13, "color": TEXT},
    {"text": "▸  100 million+ article records",          "size": 13, "color": TEXT},
    {"text": " ", "size": 7},
    {"text": "Fields used:",                             "size": 12, "color": MUTED},
    {"text": "V2Themes  ·  V2Organizations",             "size": 12, "color": MUTED, "italic": True},
    {"text": "V2Tone  ·  DocumentIdentifier",            "size": 12, "color": MUTED, "italic": True},
    {"text": " ", "size": 7},
    {"text": "Normalized per 100,000 articles",         "size": 12, "color": MUTED, "italic": True},
], 0.72, 1.98, 5.55, 5.0)

rect(s3, 6.88, 1.42, 6.1, 5.75, SURFACE, border=BORDER)
section_label(s3, "KEYWORDS TRACKED", 7.15, 1.58, GOLD)
mbox(s3, [
    {"text": "5 keyword categories",                    "size": 13, "color": MUTED},
    {"text": " ", "size": 7},
    {"text": "●  Patriotism  (爱国 / aiguo)",           "size": 15, "color": C_PAT},
    {"text": "●  China Dream  (中国梦)",                "size": 15, "color": C_CD,  "space_before": 4},
    {"text": "●  Nationalism  (民族主义)",              "size": 15, "color": C_NAT, "space_before": 4},
    {"text": "●  Guochao  (国潮)",                     "size": 15, "color": C_GC,  "space_before": 4},
    {"text": "●  Boycott  (抵制)",                     "size": 15, "color": C_BOY, "space_before": 4},
    {"text": " ", "size": 7},
    {"text": "3 findings:",                            "size": 13, "color": MUTED},
    {"text": "1  Keyword frequency over time",        "size": 13, "color": TEXT},
    {"text": "2  Early Xi vs. Consolidated Xi",       "size": 13, "color": TEXT, "space_before": 3},
    {"text": "3  Foreign vs. Domestic brand tone",   "size": 13, "color": TEXT, "space_before": 3},
], 7.15, 1.98, 5.55, 5.0)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Finding 1: Keyword Frequency (Line Chart)
# ══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
rect(s4, 0, 0, 13.33, 0.1, RED)

section_label(s4, "FINDING 1  —  KEYWORD FREQUENCY")
box(s4, "Every Keyword Grew — But Guochao Grew 56×",
    0.5, 0.68, 12.3, 0.55, size=22, bold=True, color=TEXT)
chart_subtitle(s4, "Mentions per 100k articles · 2015–2024 · Gold markers: key political events")

cd4 = ChartData()
cd4.categories = YEARS
cd4.add_series('Patriotism (爱国)',     PATRIOTISM)
cd4.add_series('China Dream (中国梦)',  CHINA_DREAM)
cd4.add_series('Nationalism (民族主义)', NATIONALISM)
cd4.add_series('Guochao (国潮)',         GUOCHAO)
cd4.add_series('Boycott (抵制)',         BOYCOTT)

gf4 = s4.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS,
    Inches(CHART_L), Inches(CHART_T), Inches(CHART_W), Inches(CHART_H), cd4)
style_line_chart(gf4.chart, [C_PAT, C_CD, C_NAT, C_GC, C_BOY])
event_markers(s4, BOYCOTT_EVENTS)

finding_callout(s4,
    "Guochao (国潮): 8 per 100k in 2015 → 496 in 2021 (+5,600%). "
    "China Dream tripled. Every keyword elevated simultaneously during boycott events.")

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Finding 2: Era Comparison (Clustered Column Chart)
# ══════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
rect(s5, 0, 0, 13.33, 0.1, RED)

section_label(s5, "FINDING 2  —  ERA COMPARISON")
box(s5, "Xi's Consolidation (2019–2024) Raised Every Baseline",
    0.5, 0.68, 12.3, 0.55, size=22, bold=True, color=TEXT)
chart_subtitle(s5, "Average mentions per 100k articles · Early Xi (2015–2018) vs. Consolidated Xi (2019–2024)")

cd5 = ChartData()
cd5.categories = ['Patriotism\n(爱国)', 'China Dream\n(中国梦)', 'Nationalism\n(民族主义)',
                  'Guochao\n(国潮)', 'Boycott\n(抵制)']
cd5.add_series('Early Xi (2015–2018)',        EARLY)
cd5.add_series('Consolidated Xi (2019–2024)', LATE)

gf5 = s5.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(CHART_L), Inches(CHART_T), Inches(CHART_W), Inches(CHART_H), cd5)
style_bar_chart(gf5.chart, [RGBColor(0xaa, 0xaa, 0xcc), RED])

finding_callout(s5,
    "Guochao: +1,685% between eras — largest jump of any keyword. "
    "Consumer nationalism grew faster than political rhetoric.")

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Finding 3: Brand Tone Divergence (Line Chart)
# ══════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
rect(s6, 0, 0, 13.33, 0.1, RED)

section_label(s6, "FINDING 3  —  BRAND TONE DIVERGENCE")
box(s6, "Buying Chinese, Boycotting Foreign: A Measurable Framing Split",
    0.5, 0.68, 12.3, 0.55, size=22, bold=True, color=TEXT)
chart_subtitle(s6,
    "GDELT tone score in Chinese state media · Negative = hostile · Positive = celebratory · 0 = neutral")

cd6 = ChartData()
cd6.categories = YEARS
cd6.add_series('Foreign Brands  (H&M, Nike, Lotte, NBA)',        FOREIGN_TONE)
cd6.add_series('Domestic Brands  (Li-Ning, Huawei, ANTA, BYD)', DOMESTIC_TONE)

gf6 = s6.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS,
    Inches(CHART_L), Inches(CHART_T), Inches(CHART_W), Inches(CHART_H), cd6)
style_line_chart(gf6.chart, [C_FOR, C_DOM])
event_markers(s6, BOYCOTT_EVENTS, boycott_only=True)

finding_callout(s6,
    "2021 divergence gap: 7.15 tone points — widest on record. "
    "Boycotts were not spontaneous public outrage. They were coordinated state media framing operations.")

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Conclusion
# ══════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
rect(s7, 0, 0, 13.33, 0.1, RED)

section_label(s7, "CONCLUSION")
box(s7, "Nationalism as a Behavioral Norm", 0.5, 0.72, 12.3, 0.65, size=30, bold=True, color=TEXT)
rect(s7, 0.5, 1.6, 12.3, 0.04, BORDER)

mbox(s7, [
    {"text": "Three measurements, one pattern:", "size": 14, "color": MUTED, "space_before": 4},
    {"text": " ", "size": 5},
    {"text": "01   Nationalist keywords grew 3–56× across a decade of state media coverage",
     "size": 15, "color": TEXT, "space_before": 8},
    {"text": "02   The 2019–2024 baseline is structurally higher — a new normal, not a spike",
     "size": 15, "color": TEXT, "space_before": 8},
    {"text": "03   Foreign brands framed as hostile · Domestic brands framed as national pride",
     "size": 15, "color": TEXT, "space_before": 8},
], 0.5, 1.75, 12.3, 2.8)

rect(s7, 0.5, 4.7, 12.3, 1.62, TEXT)
rect(s7, 0.5, 4.7, 0.08, 1.62, RED)
box(s7,
    "STATE MEDIA DID NOT REPORT ON NATIONALISM —\nIT MANUFACTURED THE CONDITIONS FOR IT.",
    0.74, 4.88, 11.8, 1.12, size=22, bold=True,
    color=RGBColor(0xff, 0xff, 0xff))

box(s7,
    "Buying a Li-Ning sneaker instead of Nike became, in this media construction, an act of patriotism.",
    0.5, 6.52, 12.3, 0.45, size=13, italic=True, color=MUTED)

# ── Save ──────────────────────────────────────────────────────────────────
OUT = "DH_Project_Presentation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  |  7 slides  |  ~5 minutes")
